import os
import json
from .models import db, Chunk, File

UPLOAD_FOLDER = os.getenv("UPLOAD_FOLDER", "data/uploads")


# -------- TEXT EXTRACTION -------- #

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            return "\n".join(page.extract_text() or "" for page in reader.pages)

        elif ext in [".txt", ".md"]:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        elif ext == ".docx":
            from docx import Document
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            return "\n".join(parts)

        elif ext in [".xlsx", ".csv"]:
            import pandas as pd
            df = pd.read_excel(file_path) if ext == ".xlsx" else pd.read_csv(file_path)
            return df.to_string(index=False)

        elif ext == ".html":
            from bs4 import BeautifulSoup
            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            return soup.get_text(separator="\n")

        elif ext == ".json":
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            return json.dumps(data, indent=2)

        else:
            return "Unsupported file type"

    except Exception as e:
        return f"Error extracting text: {e}"


# -------- TEXT CHUNKING -------- #

def chunk_text(text, chunk_size=1000, overlap=200):
    chunks = []
    start = 0
    ln = len(text)
    while start < ln:
        end = start + chunk_size
        chunks.append((text[start:end], start, end))
        start = end - overlap
    return chunks


# -------- PROCESS FILE INTO DB + FAISS -------- #

def process_file(file_id, user_id, upload_folder=None):
    if upload_folder is None:
        upload_folder = os.getenv("UPLOAD_FOLDER", "data/uploads")

    file = File.query.get(file_id)
    file_path = file.path
    created_temp = False

    if not os.path.exists(file_path):
        if file.content:
            os.makedirs(os.path.dirname(file_path), exist_ok=True)
            with open(file_path, 'wb') as f:
                f.write(file.content)
            created_temp = True
        else:
            raise FileNotFoundError(f"{file_path} does not exist and no content in DB.")

    try:
        raw_text = extract_text(file_path)
        chunked = chunk_text(raw_text)
        chunk_texts = [c[0] for c in chunked]

        from .embeddings import EmbeddingGenerator
        from .faiss_index import FaissIndex

        embedder = EmbeddingGenerator()
        faiss_index = FaissIndex(dim=embedder.get_dimension(), user_id=user_id)
        embeddings = embedder.embed_chunks(chunk_texts)

        db_chunks = []
        for text, start, end in chunked:
            chunk = Chunk(
                file_id=file.id,
                user_id=user_id,
                text=text,
                start_char=start,
                end_char=end,
                chunk_metadata={
                    "length": len(text),
                    "file_id": file.id,
                    "file_name": file.filename,
                },
            )
            db.session.add(chunk)
            db_chunks.append(chunk)

        db.session.commit()

        meta_list = []
        ids = []
        for c in db_chunks:
            meta_list.append({
                "chunk_id": c.id,
                "file_id": c.file_id,
                "file_name": file.filename,
                "start_char": c.start_char,
                "end_char": c.end_char,
                "length": len(c.text),
            })
            ids.append(c.id)

        faiss_index.add_embeddings(
            embeddings=embeddings,
            chunk_ids=ids,
            chunk_metadata=meta_list
        )

        file.processed = True
        db.session.add(file)
        db.session.commit()

    finally:
        if created_temp and os.path.exists(file_path):
            os.remove(file_path)

    return len(chunked)


# -------- FAISS SEARCH -------- #

def search_similar_chunks(user_id, query_text, top_k=5, file_ids=None):
    from .embeddings import EmbeddingGenerator
    from .faiss_index import FaissIndex

    embedder = EmbeddingGenerator()
    faiss_index = FaissIndex(dim=embedder.get_dimension(), user_id=user_id)
    query_embedding = embedder.embed_text(query_text)

    search_k = top_k * 3 if file_ids else top_k
    distances, indices, metadata = faiss_index.search(query_embedding, top_k=search_k)

    results = []
    for dist, cid, meta in zip(distances, indices, metadata):
        if cid < 0:
            continue
        cid = int(cid)
        chunk = Chunk.query.get(cid)
        if not chunk:
            faiss_index.remove_id(cid)
            continue
        if file_ids and chunk.file_id not in file_ids:
            continue
        results.append({"chunk": chunk, "distance": float(dist), "metadata": meta})
        if len(results) >= top_k:
            break

    return results
